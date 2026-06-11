"""PowerPoint (.pptx) parser with section_ref support."""

from pathlib import Path

from .base import BaseParser, ParseResult, Section


class PPTXParser(BaseParser):
    """Parser for .pptx files using python-pptx."""

    def parse(self, file_path: Path) -> ParseResult:
        from pptx import Presentation

        prs = Presentation(str(file_path))
        sections: list[Section] = []
        all_text_parts: list[str] = []

        for idx, slide in enumerate(prs.slides, start=1):
            # Extract slide title
            title = f"Slide {idx}"
            if slide.shapes.title and slide.shapes.title.text.strip():
                title = slide.shapes.title.text.strip()

            # Extract all text from shapes
            text_parts: list[str] = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for paragraph in shape.text_frame.paragraphs:
                        line = paragraph.text.strip()
                        if line:
                            text_parts.append(line)

            # Extract speaker notes
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                notes = slide.notes_slide.notes_text_frame.text.strip()
                if notes:
                    text_parts.append(f"\n[Speaker Notes]\n{notes}")

            content = "\n".join(text_parts)
            sections.append(Section(
                title=title,
                content=content,
                section_ref=f"slide-{idx}",
                page_start=idx,
            ))
            all_text_parts.append(f"=== {title} ===\n{content}")

        slide_count = len(prs.slides)

        if not sections:
            sections = [Section(
                title="Empty Presentation",
                content="No slides found.",
                section_ref="slide-0",
            )]

        return ParseResult(
            filename=file_path.name,
            sections=sections,
            raw_text="\n\n".join(all_text_parts),
            page_count=slide_count,
            metadata={
                "parser": "python-pptx",
                "slide_count": slide_count,
            },
        )
